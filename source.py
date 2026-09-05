import os
import io
import sys
import threading
import urllib.request
import json
import locale
import tkinter as tk
from tkinter import filedialog, messagebox, colorchooser
import customtkinter as ctk
import webbrowser
from PIL import Image, ImageDraw, ImageOps
import pymupdf
from moviepy import VideoFileClip, AudioFileClip
import yt_dlp

if sys.platform == "win32":
    import winreg

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

ctk.set_appearance_mode("Light")
ctk.set_default_color_theme("blue")

BG_SIDEBAR = "#F1F3F5"
BG_MAIN = "#FFFFFF"
BG_CARD = "#F8F9FA"
TEXT_MAIN = "#1A1A1A"
TEXT_MUTED = "#707070"

CURRENT_VERSION = "1.0"
GITHUB_REPO = "Luv-Kitty/LapisTool" 
REG_KEY_PATH = r"Software\luv_kitty_Converter"

TRANSLATIONS = {
    "ru": {
        "app_title": "luv_kitty's Converter — Конвертер & Загрузчик",
        "nav_conv": " 🔄  Конвертер",
        "nav_down": " 📥  Загрузчик",
        "nav_about": " 👤  Об авторе",
        "nav_settings": " ⚙  Настройки",
        "version": "Версия",
        "conv_title": "Конвертер файлов",
        "btn_select_file": "Выбрать файл",
        "no_file": "Файл не выбран",
        "format_lbl": "Формат:",
        "codec_lbl": "Кодек:",
        "waiting_file": "Ожидание файла...",
        "btn_convert": "Начать конвертацию",
        "down_title": "Загрузка медиа по ссылке",
        "down_sub": "Поддерживаемые сервисы: YouTube, VK Видео, Rutube, Mail.ru и др.",
        "url_placeholder": "Вставьте ссылку на видео...",
        "folder_btn": "Папка",
        "folder_lbl": "Папка:",
        "ready_dl": "Готов к скачиванию",
        "btn_download": "Скачать медиа",
        "about_title": "Об авторе и проекте",
        "about_desc": "luv_kitty's Converter — компактный инструмент для конвертации\nмедиафайлов и удобного скачивания роликов из сети.",
        "github_btn": "GitHub Профиль",
        "settings_title": "Настройки программы",
        "lang_lbl": "Язык интерфейса (Language):",
        "color_lbl": "Основной цвет (RGB):",
        "color_btn": "Выбрать цвет",
        "check_update_btn": "Проверить обновления",
        "update_avail": "Доступно обновление на GitHub! Версия:",
        "update_no": "У вас установлена последняя версия.",
        "update_title": "Обновление",
        "update_error": "Ошибка при проверке обновлений",
        "select_file_dialog": "Все поддерживаемые",
        "img_dialog": "Изображения",
        "vid_dialog": "Видео",
        "aud_dialog": "Аудио",
        "doc_dialog": "Документы и Презентации",
        "ready_process": "Готов к обработке",
        "converting": "Конвертируем...",
        "success_saved": "Успешно сохранено!",
        "done": "Готово",
        "file_converted": "Файл успешно сконвертирован!",
        "error": "Ошибка",
        "fail": "Сбой:",
        "warning": "Внимание",
        "enter_url": "Пожалуйста, введите ссылку!",
        "connecting": "Подключение к сервису...",
        "downloading": "Загрузка файла",
        "dl_complete": "Загрузка завершена!",
        "success": "Успех",
        "media_downloaded": "Медиафайл успешно загружен!",
        "dl_error": "Ошибка загрузки",
        "dl_failed": "Не удалось скачать видео:"
    },
    "en": {
        "app_title": "luv_kitty's Converter — Converter & Downloader",
        "nav_conv": " 🔄  Converter",
        "nav_down": " 📥  Downloader",
        "nav_about": " 👤  About",
        "nav_settings": " ⚙  Settings",
        "version": "Version",
        "conv_title": "File Converter",
        "btn_select_file": "Select File",
        "no_file": "No file selected",
        "format_lbl": "Format:",
        "codec_lbl": "Codec:",
        "waiting_file": "Waiting for file...",
        "btn_convert": "Start Conversion",
        "down_title": "Download Media by URL",
        "down_sub": "Supported: YouTube, VK Video, Rutube, etc.",
        "url_placeholder": "Paste video URL here...",
        "folder_btn": "Folder",
        "folder_lbl": "Folder:",
        "ready_dl": "Ready to download",
        "btn_download": "Download Media",
        "about_title": "About the Author & Project",
        "about_desc": "luv_kitty's Converter is a compact tool for converting\nmedia files and easily downloading videos from the web.",
        "github_btn": "GitHub Profile",
        "settings_title": "Application Settings",
        "lang_lbl": "Interface Language:",
        "color_lbl": "Primary Color (RGB):",
        "color_btn": "Pick Color",
        "check_update_btn": "Check for Updates",
        "update_avail": "An update is available on GitHub! Version:",
        "update_no": "You are using the latest version.",
        "update_title": "Update Check",
        "update_error": "Error checking for updates",
        "select_file_dialog": "All supported",
        "img_dialog": "Images",
        "vid_dialog": "Video",
        "aud_dialog": "Audio",
        "doc_dialog": "Documents & Presentations",
        "ready_process": "Ready to process",
        "converting": "Converting...",
        "success_saved": "Successfully saved!",
        "done": "Done",
        "file_converted": "File converted successfully!",
        "error": "Error",
        "fail": "Failure:",
        "warning": "Warning",
        "enter_url": "Please enter a URL!",
        "connecting": "Connecting to service...",
        "downloading": "Downloading file",
        "dl_complete": "Download completed!",
        "success": "Success",
        "media_downloaded": "Media file successfully downloaded!",
        "dl_error": "Download Error",
        "dl_failed": "Failed to download video:"
    }
}

class luvkittysApp(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.red_primary = "#E53935"
        self.red_hover = "#C62828"
        self.download_folder = os.path.expanduser("~/Downloads")
        
        try:
            sys_lang = locale.getdefaultlocale()[0]
            self.current_lang = 'ru' if sys_lang and sys_lang.startswith('ru') else 'en'
        except:
            self.current_lang = 'en'

        self.saved_conv_codec = None
        self.saved_dl_type = None
        self.saved_dl_codec = None

        self.load_settings()

        self.geometry("900x520")
        self.resizable(False, False)
        self.configure(fg_color=BG_MAIN)

        self.filepath = ""
        self.file_type = None
        self.current_tab = None
        self._anim_id = None
        self.avatar_image = None

        self.setup_ui()
        self.apply_language()
        self.apply_colors()
        self.restore_saved_choices()
        
        threading.Thread(target=self.load_github_avatar, daemon=True).start()
        threading.Thread(target=self.check_for_updates, args=(False,), daemon=True).start()

    def load_settings(self):
        if sys.platform == "win32":
            try:
                key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, REG_KEY_PATH, 0, winreg.KEY_READ)
                self.current_lang, _ = winreg.QueryValueEx(key, "lang")
                self.red_primary, _ = winreg.QueryValueEx(key, "color")
                self.download_folder, _ = winreg.QueryValueEx(key, "download_folder")
                self.saved_conv_codec, _ = winreg.QueryValueEx(key, "conv_codec")
                self.saved_dl_type, _ = winreg.QueryValueEx(key, "dl_type")
                self.saved_dl_codec, _ = winreg.QueryValueEx(key, "dl_codec")
                winreg.CloseKey(key)
                self._update_hover_color()
            except Exception:
                pass

    def save_settings(self, *args):
        if sys.platform == "win32":
            try:
                key = winreg.CreateKey(winreg.HKEY_CURRENT_USER, REG_KEY_PATH)
                winreg.SetValueEx(key, "lang", 0, winreg.REG_SZ, str(self.current_lang))
                winreg.SetValueEx(key, "color", 0, winreg.REG_SZ, str(self.red_primary))
                winreg.SetValueEx(key, "download_folder", 0, winreg.REG_SZ, str(self.download_folder))
                
                conv_codec = self.combo_codec.get() if hasattr(self, 'combo_codec') else ""
                dl_type = self.dl_type_combo.get() if hasattr(self, 'dl_type_combo') else ""
                dl_codec = self.dl_codec_combo.get() if hasattr(self, 'dl_codec_combo') else ""

                winreg.SetValueEx(key, "conv_codec", 0, winreg.REG_SZ, str(conv_codec))
                winreg.SetValueEx(key, "dl_type", 0, winreg.REG_SZ, str(dl_type))
                winreg.SetValueEx(key, "dl_codec", 0, winreg.REG_SZ, str(dl_codec))
                winreg.CloseKey(key)
            except Exception as e:
                print(f"Ошибка сохранения настроек в реестр: {e}")

    def restore_saved_choices(self):
        if self.saved_conv_codec and self.saved_conv_codec in self.combo_codec._values:
            self.combo_codec.set(self.saved_conv_codec)
        if self.saved_dl_type and self.saved_dl_type in self.dl_type_combo._values:
            self.dl_type_combo.set(self.saved_dl_type)
        if self.saved_dl_codec and self.saved_dl_codec in self.dl_codec_combo._values:
            self.dl_codec_combo.set(self.saved_dl_codec)

    def _update_hover_color(self):
        hex_color = self.red_primary.lstrip('#')
        try:
            r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            self.red_hover = f"#{max(0, r - 30):02x}{max(0, g - 30):02x}{max(0, b - 30):02x}"
        except:
            self.red_hover = self.red_primary

    def _(self, key):
        return TRANSLATIONS.get(self.current_lang, TRANSLATIONS["en"]).get(key, key)

    def setup_ui(self):
        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)

        self.sidebar = ctk.CTkFrame(self, width=220, fg_color=BG_SIDEBAR, corner_radius=0)
        self.sidebar.grid(row=0, column=0, sticky="nsew")
        self.sidebar.grid_rowconfigure(5, weight=1)

        self.logo_label = ctk.CTkLabel(
            self.sidebar,
            text="luv_kitty's",
            font=ctk.CTkFont(size=24, weight="bold"),
            text_color=self.red_primary,
        )
        self.logo_label.pack(pady=(25, 2), padx=20, anchor="w")

        self.sub_logo = ctk.CTkLabel(
            self.sidebar,
            text="MEDIA SUITE",
            font=ctk.CTkFont(size=11, weight="bold"),
            text_color=TEXT_MUTED,
        )
        self.sub_logo.pack(pady=(0, 25), padx=20, anchor="w")

        self.btn_nav_conv = ctk.CTkButton(
            self.sidebar, text="", anchor="w", font=ctk.CTkFont(size=14, weight="bold"),
            fg_color="transparent", text_color=TEXT_MAIN, hover_color="#DCE1E5", height=42, corner_radius=8,
            command=lambda: self.select_tab("converter")
        )
        self.btn_nav_conv.pack(fill="x", padx=15, pady=5)

        self.btn_nav_down = ctk.CTkButton(
            self.sidebar, text="", anchor="w", font=ctk.CTkFont(size=14, weight="bold"),
            fg_color="transparent", text_color=TEXT_MAIN, hover_color="#DCE1E5", height=42, corner_radius=8,
            command=lambda: self.select_tab("downloader")
        )
        self.btn_nav_down.pack(fill="x", padx=15, pady=5)

        self.btn_nav_about = ctk.CTkButton(
            self.sidebar, text="", anchor="w", font=ctk.CTkFont(size=14, weight="bold"),
            fg_color="transparent", text_color=TEXT_MAIN, hover_color="#DCE1E5", height=42, corner_radius=8,
            command=lambda: self.select_tab("about")
        )
        self.btn_nav_about.pack(fill="x", padx=15, pady=5)
        
        self.btn_nav_settings = ctk.CTkButton(
            self.sidebar, text="", anchor="w", font=ctk.CTkFont(size=14, weight="bold"),
            fg_color="transparent", text_color=TEXT_MAIN, hover_color="#DCE1E5", height=42, corner_radius=8,
            command=lambda: self.select_tab("settings")
        )
        self.btn_nav_settings.pack(fill="x", padx=15, pady=5)

        self.lbl_version = ctk.CTkLabel(self.sidebar, text="", font=ctk.CTkFont(size=11), text_color=TEXT_MUTED)
        self.lbl_version.pack(side="bottom", pady=15, padx=20, anchor="w")

        self.main_container = ctk.CTkFrame(self, fg_color=BG_MAIN, corner_radius=0)
        self.main_container.grid(row=0, column=1, sticky="nsew", padx=25, pady=20)

        self.tab_frames = {}
        self.tab_frames["converter"] = self.build_converter_view()
        self.tab_frames["downloader"] = self.build_downloader_view()
        self.tab_frames["about"] = self.build_about_view()
        self.tab_frames["settings"] = self.build_settings_view()

        self.select_tab("converter")

    def apply_language(self):
        self.title(self._("app_title"))
        self.btn_nav_conv.configure(text=self._("nav_conv"))
        self.btn_nav_down.configure(text=self._("nav_down"))
        self.btn_nav_about.configure(text=self._("nav_about"))
        self.btn_nav_settings.configure(text=self._("nav_settings"))
        self.lbl_version.configure(text=f"{self._('version')} {CURRENT_VERSION}")

        vid = "Видео" if self.current_lang == "ru" else "Video"
        aud = "Аудио" if self.current_lang == "ru" else "Audio"
        page = "страницы" if self.current_lang == "ru" else "pages"
        txt = "текст" if self.current_lang == "ru" else "text"

        self.formats = {
            "image": ["PNG", "JPG", "WEBP", "BMP", "ICO", "PDF", "TIFF"],
            "video": ["MP4", "AVI", "MKV", "MOV", "GIF", "WEBM", f"MP3 ({aud})", f"WAV ({aud})", f"FLAC ({aud})"],
            "audio": ["MP3", "WAV", "OGG", "M4A", "FLAC", "AAC"],
            "presentation": ["PDF", f"PNG ({page})", f"JPG ({page})", f"TXT ({txt})"],
        }

        self.conv_title.configure(text=self._("conv_title"))
        self.btn_select.configure(text=self._("btn_select_file"))
        if not self.filepath:
            self.lbl_filename.configure(text=self._("no_file"))
        self.lbl_target.configure(text=self._("format_lbl"))
        self.lbl_codec.configure(text=self._("codec_lbl"))
        if self.progress_conv.get() == 0:
            self.lbl_status_conv.configure(text=self._("waiting_file"))
        self.btn_convert.configure(text=self._("btn_convert"))
        if self.file_type:
            opts = self.formats.get(self.file_type, [])
            if opts:
                self.combo_format.configure(values=opts)

        self.down_title.configure(text=self._("down_title"))
        self.lbl_sub.configure(text=self._("down_sub"))
        self.url_entry.configure(placeholder_text=self._("url_placeholder"))
        self.lbl_fmt.configure(text=self._("format_lbl"))
        self.lbl_dl_codec.configure(text=self._("codec_lbl"))
        self.btn_folder.configure(text=self._("folder_btn"))
        self.lbl_folder_path.configure(text=f"{self._('folder_lbl')} {self.download_folder}")
        if self.progress_dl.get() == 0:
            self.lbl_status_dl.configure(text=self._("ready_dl"))
        self.btn_download.configure(text=self._("btn_download"))
        
        current_dl = self.dl_type_combo.get()
        new_dl_values = [
            f"MP4 ({vid})", f"WEBM ({vid})", f"MKV ({vid})",
            f"MP3 ({aud})", f"WAV ({aud})", f"FLAC ({aud})", f"M4A ({aud})"
        ]
        self.dl_type_combo.configure(values=new_dl_values)
        if current_dl == "" or current_dl.startswith("MP4"):
            self.dl_type_combo.set(new_dl_values[0])

        self.about_title.configure(text=self._("about_title"))
        self.lbl_desc.configure(text=self._("about_desc"))
        self.btn_github.configure(text=self._("github_btn"))

        self.settings_title.configure(text=self._("settings_title"))
        self.lbl_lang.configure(text=self._("lang_lbl"))
        self.lbl_color.configure(text=self._("color_lbl"))
        self.btn_color.configure(text=self._("color_btn"))
        self.btn_check_update.configure(text=self._("check_update_btn"))
        self.combo_lang.set("Русский" if self.current_lang == "ru" else "English")

    def apply_colors(self):
        self.logo_label.configure(text_color=self.red_primary)
        self.select_tab(self.current_tab, force_color=True)

        self.btn_select.configure(fg_color=self.red_primary, hover_color=self.red_hover)
        self.combo_format.configure(button_color=self.red_primary, button_hover_color=self.red_hover)
        self.combo_codec.configure(button_color=self.red_primary, button_hover_color=self.red_hover)
        self.progress_conv.configure(progress_color=self.red_primary)
        self.btn_convert.configure(fg_color=self.red_primary, hover_color=self.red_hover)
        
        self.url_entry.configure(border_color=self.red_primary)
        self.dl_type_combo.configure(button_color=self.red_primary, button_hover_color=self.red_hover)
        self.dl_codec_combo.configure(button_color=self.red_primary, button_hover_color=self.red_hover)
        self.progress_dl.configure(progress_color=self.red_primary)
        self.btn_download.configure(fg_color=self.red_primary, hover_color=self.red_hover)
        
        self.combo_lang.configure(button_color=self.red_primary, button_hover_color=self.red_hover)
        self.btn_color.configure(fg_color=self.red_primary, hover_color=self.red_hover)
        self.btn_check_update.configure(fg_color=self.red_primary, hover_color=self.red_hover)

    def change_color(self):
        color = colorchooser.askcolor(title=self._("color_btn"))
        if color[1]:
            self.red_primary = color[1]
            self._update_hover_color()
            self.apply_colors()
            self.save_settings()

    def change_language(self, choice):
        self.current_lang = "ru" if choice == "Русский" else "en"
        self.apply_language()
        self.save_settings()

    def check_for_updates(self, manual=False):
        try:
            url = f"https://api.github.com/repos/{GITHUB_REPO}/releases/latest"
            req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req, timeout=4) as response:
                data = json.loads(response.read().decode())
                latest = data.get("tag_name", "").replace("v", "")
                
                if latest and latest != CURRENT_VERSION:
                    self.after(0, lambda: self.show_custom_update_dialog(latest))
                elif manual:
                    self.after(0, lambda: messagebox.showinfo(self._("update_title"), self._('update_no')))
        except Exception as e:
            if manual:
                self.after(0, lambda: messagebox.showerror(self._("update_title"), f"{self._('update_error')}:\n{str(e)}"))

    def show_custom_update_dialog(self, latest_version):
        dialog = ctk.CTkToplevel(self)
        dialog.title(self._("update_title"))
        dialog.geometry("400x220")
        dialog.resizable(False, False)
        dialog.attributes("-topmost", True)
        dialog.grab_set()

        x = self.winfo_x() + (self.winfo_width() // 2) - 200
        y = self.winfo_y() + (self.winfo_height() // 2) - 110
        dialog.geometry(f"+{x}+{y}")
        
        lbl_title_text = "Доступно обновление! 🎉" if self.current_lang == "ru" else "Update available! 🎉"
        lbl_info_text = f"Доступна новая версия {latest_version}.\nУ вас установлена {CURRENT_VERSION}." if self.current_lang == "ru" else f"New version {latest_version} is available.\nYou have {CURRENT_VERSION}."
        btn_dl_text = "Скачать" if self.current_lang == "ru" else "Download"
        btn_skip_text = "Пропустить" if self.current_lang == "ru" else "Skip"

        lbl_title = ctk.CTkLabel(dialog, text=lbl_title_text, font=ctk.CTkFont(size=18, weight="bold"))
        lbl_title.pack(pady=(20, 5))

        lbl_info = ctk.CTkLabel(dialog, text=lbl_info_text, text_color=TEXT_MUTED)
        lbl_info.pack(pady=(0, 20))

        btn_frame = ctk.CTkFrame(dialog, fg_color="transparent")
        btn_frame.pack(fill="x", padx=20)

        btn_download = ctk.CTkButton(btn_frame, text=btn_dl_text, fg_color=self.red_primary, hover_color=self.red_hover, 
                                     command=lambda: [webbrowser.open(f"https://github.com/{GITHUB_REPO}/releases/latest"), dialog.destroy()])
        btn_download.pack(side="left", expand=True, padx=5)

        btn_skip = ctk.CTkButton(btn_frame, text=btn_skip_text, fg_color="transparent", border_width=1, 
                                 text_color=TEXT_MAIN, command=dialog.destroy)
        btn_skip.pack(side="right", expand=True, padx=5)

    def select_tab(self, name, force_color=False):
        if self.current_tab == name and not force_color:
            return

        if self._anim_id:
            self.after_cancel(self._anim_id)
            self._anim_id = None

        self.current_tab = name

        buttons = [("converter", self.btn_nav_conv), ("downloader", self.btn_nav_down), 
                   ("about", self.btn_nav_about), ("settings", self.btn_nav_settings)]
        for tab_name, btn in buttons:
            if tab_name == name:
                btn.configure(fg_color=self.red_primary, text_color="white", hover_color=self.red_hover)
            else:
                btn.configure(fg_color="transparent", text_color=TEXT_MAIN, hover_color="#DCE1E5")

        if not force_color:
            for frame in self.tab_frames.values():
                frame.place_forget()
            target_frame = self.tab_frames[name]
            self._animate_tab_in(target_frame, step=0)

    def _animate_tab_in(self, frame, step=0):
        total_steps = 15
        if step <= total_steps:
            t = step / total_steps
            ease_out = 1 - pow(1 - t, 3)
            offset = 0.05 * (1 - ease_out)
            frame.place(relx=0, rely=offset, relwidth=1, relheight=1)
            self._anim_id = self.after(10, lambda: self._animate_tab_in(frame, step + 1))
        else:
            frame.place(relx=0, rely=0, relwidth=1, relheight=1)
            self._anim_id = None

    def load_github_avatar(self):
        size = (90, 90)
        try:
            url = "https://github.com/Luv-Kitty.png"
            req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req, timeout=4) as resp:
                raw_data = resp.read()

            img = Image.open(io.BytesIO(raw_data)).convert("RGBA")
            mask = Image.new('L', (size[0] * 4, size[1] * 4), 0)
            draw = ImageDraw.Draw(mask)
            draw.ellipse((0, 0, size[0] * 4, size[1] * 4), fill=255)
            mask = mask.resize(size, Image.Resampling.LANCZOS)

            output = ImageOps.fit(img, size, centering=(0.5, 0.5))
            circle_img = Image.new("RGBA", size, (0, 0, 0, 0))
            circle_img.paste(output, (0, 0), mask=mask)
            self.avatar_image = ctk.CTkImage(light_image=circle_img, dark_image=circle_img, size=size)
        except Exception:
            fallback = Image.new("RGBA", size, (0, 0, 0, 0))
            draw = ImageDraw.Draw(fallback)
            draw.ellipse((0, 0, 89, 89), fill=self.red_primary)
            self.avatar_image = ctk.CTkImage(light_image=fallback, dark_image=fallback, size=size)

        self.after(0, self._apply_avatar)

    def _apply_avatar(self):
        if hasattr(self, 'lbl_avatar_widget') and self.avatar_image:
            self.lbl_avatar_widget.configure(image=self.avatar_image, text="")

    def build_converter_view(self):
        frame = ctk.CTkFrame(self.main_container, fg_color="transparent")
        self.conv_title = ctk.CTkLabel(frame, text="", font=ctk.CTkFont(size=20, weight="bold"), text_color=TEXT_MAIN)
        self.conv_title.pack(anchor="w", pady=(0, 15))

        card = ctk.CTkFrame(frame, fg_color=BG_CARD, corner_radius=14)
        card.pack(fill="both", expand=True)

        self.btn_select = ctk.CTkButton(
            card, text="", command=self.select_file, height=40, width=180,
            font=ctk.CTkFont(size=14, weight="bold"), corner_radius=8
        )
        self.btn_select.pack(pady=(20, 8))

        self.lbl_filename = ctk.CTkLabel(card, text="", text_color=TEXT_MUTED, font=ctk.CTkFont(size=13))
        self.lbl_filename.pack(pady=(0, 15))

        fmt_frame = ctk.CTkFrame(card, fg_color="transparent")
        fmt_frame.pack(pady=5)

        self.lbl_target = ctk.CTkLabel(fmt_frame, text="", text_color=TEXT_MAIN, font=ctk.CTkFont(size=13, weight="bold"))
        self.lbl_target.pack(side="left", padx=(0, 5))

        self.combo_format = ctk.CTkComboBox(fmt_frame, values=[""], state="disabled", width=140, height=36)
        self.combo_format.pack(side="left", padx=(0, 15))

        self.lbl_codec = ctk.CTkLabel(fmt_frame, text="", text_color=TEXT_MAIN, font=ctk.CTkFont(size=13, weight="bold"))
        self.lbl_codec.pack(side="left", padx=(0, 5))

        codecs_list = [
            "H.264 (libx264)",
            "H.265 (libx265)",
            "NVENC H.264 (h264_nvenc)",
            "NVENC HEVC (hevc_nvenc)",
            "AMD H.264 (h264_amf)",
            "Intel QSV (h264_qsv)",
            "VP9 (libvpx-vp9)",
            "VP8 (libvpx)",
            "AV1 (libsvtav1)",
            "MPEG4",
            "ProRes (prores_ks)",
            "DNxHD"
        ]

        self.combo_codec = ctk.CTkComboBox(
            fmt_frame, values=codecs_list, state="disabled", width=190, height=36, command=self.save_settings
        )
        self.combo_codec.pack(side="left")

        self.progress_conv = ctk.CTkProgressBar(card, width=480, height=10, corner_radius=5)
        self.progress_conv.pack(pady=(25, 8))
        self.progress_conv.set(0)

        self.lbl_status_conv = ctk.CTkLabel(card, text="", text_color=TEXT_MUTED, font=ctk.CTkFont(size=12))
        self.lbl_status_conv.pack(pady=(0, 15))

        self.btn_convert = ctk.CTkButton(
            card, text="", command=self.start_conversion_thread,
            text_color="white", text_color_disabled="#EEEEEE", height=44, width=220,
            font=ctk.CTkFont(size=14, weight="bold"), corner_radius=8, state="disabled"
        )
        self.btn_convert.pack(pady=(0, 20))
        return frame

    def build_downloader_view(self):
        frame = ctk.CTkFrame(self.main_container, fg_color="transparent")
        self.down_title = ctk.CTkLabel(frame, text="", font=ctk.CTkFont(size=20, weight="bold"), text_color=TEXT_MAIN)
        self.down_title.pack(anchor="w", pady=(0, 15))

        card = ctk.CTkFrame(frame, fg_color=BG_CARD, corner_radius=14)
        card.pack(fill="both", expand=True)

        self.lbl_sub = ctk.CTkLabel(card, text="", text_color=TEXT_MUTED, font=ctk.CTkFont(size=12))
        self.lbl_sub.pack(pady=(15, 10))

        self.url_entry = ctk.CTkEntry(card, placeholder_text="", width=540, height=42, corner_radius=8)
        self.url_entry.pack(pady=5)

        def paste_handler(event=None):
            try:
                self.url_entry.insert("insert", self.clipboard_get())
                return "break"
            except tk.TclError:
                pass

        for key in ["<Control-v>", "<Control-V>", "<Control-м>", "<Control-М>", "<Command-v>"]:
            self.url_entry.bind(key, paste_handler)

        opts_row = ctk.CTkFrame(card, fg_color="transparent")
        opts_row.pack(pady=12)

        self.lbl_fmt = ctk.CTkLabel(opts_row, text="", text_color=TEXT_MAIN, font=ctk.CTkFont(size=13, weight="bold"))
        self.lbl_fmt.pack(side="left", padx=(0, 5))

        self.dl_type_combo = ctk.CTkComboBox(opts_row, values=[""], width=140, height=34, command=self.save_settings)
        self.dl_type_combo.pack(side="left", padx=(0, 10))

        self.lbl_dl_codec = ctk.CTkLabel(opts_row, text="", text_color=TEXT_MAIN, font=ctk.CTkFont(size=13, weight="bold"))
        self.lbl_dl_codec.pack(side="left", padx=(0, 5))

        dl_codecs = ["Auto", "H.264 (AVC)", "H.265 (HEVC)", "VP9", "AV1", "ProRes", "DNxHD"]
        self.dl_codec_combo = ctk.CTkComboBox(
            opts_row, values=dl_codecs, width=120, height=34, command=self.save_settings
        )
        self.dl_codec_combo.pack(side="left", padx=(0, 10))

        self.btn_folder = ctk.CTkButton(
            opts_row, text="", command=self.choose_download_folder,
            fg_color="#4B5563", hover_color="#374151", text_color="white", height=34, width=90, corner_radius=6
        )
        self.btn_folder.pack(side="left")

        self.lbl_folder_path = ctk.CTkLabel(card, text="", text_color=TEXT_MUTED, font=ctk.CTkFont(size=11))
        self.lbl_folder_path.pack(pady=(0, 10))

        self.progress_dl = ctk.CTkProgressBar(card, width=540, height=10, corner_radius=5)
        self.progress_dl.pack(pady=(10, 8))
        self.progress_dl.set(0)

        self.lbl_status_dl = ctk.CTkLabel(card, text="", text_color=TEXT_MUTED, font=ctk.CTkFont(size=12))
        self.lbl_status_dl.pack(pady=(0, 12))

        self.btn_download = ctk.CTkButton(
            card, text="", command=self.start_download_thread, text_color="white",
            height=44, width=220, font=ctk.CTkFont(size=14, weight="bold"), corner_radius=8
        )
        self.btn_download.pack(pady=(0, 15))
        return frame

    def build_about_view(self):
        frame = ctk.CTkFrame(self.main_container, fg_color="transparent")
        self.about_title = ctk.CTkLabel(frame, text="", font=ctk.CTkFont(size=20, weight="bold"), text_color=TEXT_MAIN)
        self.about_title.pack(anchor="w", pady=(0, 15))

        card = ctk.CTkFrame(frame, fg_color=BG_CARD, corner_radius=14)
        card.pack(fill="both", expand=True)

        inner_layout = ctk.CTkFrame(card, fg_color="transparent")
        inner_layout.pack(expand=True, pady=30, padx=30)

        self.lbl_avatar_widget = ctk.CTkLabel(inner_layout, text="", width=90, height=90)
        self.lbl_avatar_widget.pack(side="left", padx=(0, 25))

        info_box = ctk.CTkFrame(inner_layout, fg_color="transparent")
        info_box.pack(side="left", fill="both", expand=True)

        self.lbl_author = ctk.CTkLabel(info_box, text="luv_kitty", font=ctk.CTkFont(size=22, weight="bold"), text_color=TEXT_MAIN)
        self.lbl_author.pack(anchor="w", pady=(0, 5))

        self.lbl_desc = ctk.CTkLabel(info_box, text="", font=ctk.CTkFont(size=13), text_color=TEXT_MUTED, justify="left")
        self.lbl_desc.pack(anchor="w", pady=(0, 15))

        github_url = "https://github.com/Luv-Kitty"
        self.btn_github = ctk.CTkButton(
            info_box, text="", command=lambda: webbrowser.open(github_url),
            fg_color="#1F2937", hover_color="#111827", text_color="white",
            height=38, width=160, font=ctk.CTkFont(size=13, weight="bold"), corner_radius=8
        )
        self.btn_github.pack(anchor="w")
        return frame

    def build_settings_view(self):
        frame = ctk.CTkFrame(self.main_container, fg_color="transparent")
        self.settings_title = ctk.CTkLabel(frame, text="", font=ctk.CTkFont(size=20, weight="bold"), text_color=TEXT_MAIN)
        self.settings_title.pack(anchor="w", pady=(0, 15))

        card = ctk.CTkFrame(frame, fg_color=BG_CARD, corner_radius=14)
        card.pack(fill="both", expand=True)

        lang_row = ctk.CTkFrame(card, fg_color="transparent")
        lang_row.pack(fill="x", padx=30, pady=(30, 10))

        self.lbl_lang = ctk.CTkLabel(lang_row, text="", text_color=TEXT_MAIN, font=ctk.CTkFont(size=14, weight="bold"))
        self.lbl_lang.pack(side="left")

        self.combo_lang = ctk.CTkComboBox(lang_row, values=["English", "Русский"], command=self.change_language)
        self.combo_lang.pack(side="right")

        color_row = ctk.CTkFrame(card, fg_color="transparent")
        color_row.pack(fill="x", padx=30, pady=10)

        self.lbl_color = ctk.CTkLabel(color_row, text="", text_color=TEXT_MAIN, font=ctk.CTkFont(size=14, weight="bold"))
        self.lbl_color.pack(side="left")

        self.btn_color = ctk.CTkButton(
            color_row, text="", command=self.change_color, text_color="white", font=ctk.CTkFont(size=13, weight="bold")
        )
        self.btn_color.pack(side="right")

        update_row = ctk.CTkFrame(card, fg_color="transparent")
        update_row.pack(fill="x", padx=30, pady=20)

        self.btn_check_update = ctk.CTkButton(
            update_row, text="", command=lambda: threading.Thread(target=self.check_for_updates, args=(True,), daemon=True).start(),
            text_color="white", height=40, font=ctk.CTkFont(size=14, weight="bold")
        )
        self.btn_check_update.pack(fill="x")
        return frame

    def select_file(self):
        path = filedialog.askopenfilename(
            filetypes=[
                (self._("select_file_dialog"), "*.png *.jpg *.jpeg *.webp *.bmp *.tiff *.mp4 *.avi *.mkv *.mov *.webm *.mp3 *.wav *.ogg *.m4a *.flac *.aac *.pdf *.pptx"),
                (self._("img_dialog"), "*.png *.jpg *.jpeg *.webp *.bmp *.tiff"),
                (self._("vid_dialog"), "*.mp4 *.avi *.mkv *.mov *.webm"),
                (self._("aud_dialog"), "*.mp3 *.wav *.ogg *.m4a *.flac *.aac"),
                (self._("doc_dialog"), "*.pdf *.pptx"),
            ]
        )
        if not path:
            return

        self.filepath = path
        ext = os.path.splitext(path)[1].lower()
        self.lbl_filename.configure(text=os.path.basename(path), text_color=TEXT_MAIN)

        if ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"]:
            self.file_type = "image"
        elif ext in [".mp4", ".avi", ".mkv", ".mov", ".webm"]:
            self.file_type = "video"
        elif ext in [".mp3", ".wav", ".ogg", ".m4a", ".flac", ".aac"]:
            self.file_type = "audio"
        elif ext in [".pdf", ".pptx"]:
            self.file_type = "presentation"

        opts = self.formats.get(self.file_type, [])
        if opts:
            self.combo_format.configure(values=opts, state="normal")
            self.combo_format.set(opts[0])
            self.combo_codec.configure(state="normal" if self.file_type == "video" else "disabled")
            self.btn_convert.configure(state="normal")
            self.lbl_status_conv.configure(text=f"{self._('ready_process')} ({self.file_type.capitalize()})")

    def start_conversion_thread(self):
        self.btn_convert.configure(state="disabled")
        self.progress_conv.set(0.3)
        self.lbl_status_conv.configure(text=self._("converting"))
        threading.Thread(target=self.run_conversion, daemon=True).start()

    def run_conversion(self):
        raw_target = self.combo_format.get()
        target_fmt = raw_target.split()[0].lower()
        out_dir = os.path.dirname(self.filepath)
        base_name = os.path.splitext(os.path.basename(self.filepath))[0]

        codec_map = {
            "H.264 (libx264)": "libx264", 
            "H.265 (libx265)": "libx265", 
            "NVENC H.264 (h264_nvenc)": "h264_nvenc",
            "NVENC HEVC (hevc_nvenc)": "hevc_nvenc",
            "AMD H.264 (h264_amf)": "h264_amf",
            "Intel QSV (h264_qsv)": "h264_qsv",
            "VP9 (libvpx-vp9)": "libvpx-vp9",
            "VP8 (libvpx)": "libvpx",
            "AV1 (libsvtav1)": "libsvtav1",
            "MPEG4": "mpeg4",
            "ProRes (prores_ks)": "prores_ks",
            "DNxHD": "dnxhd"
        }
        selected_codec = codec_map.get(self.combo_codec.get(), "libx264")

        try:
            if self.file_type == "image":
                img = Image.open(self.filepath)
                if target_fmt in ["jpg", "jpeg", "bmp"] and img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                save_path = os.path.join(out_dir, f"{base_name}.{target_fmt}")
                if target_fmt == "ico":
                    img.resize((256, 256)).save(save_path, format="ICO")
                elif target_fmt == "pdf":
                    img.convert("RGB").save(save_path, "PDF")
                else:
                    img.save(save_path)

            elif self.file_type == "video":
                clip = VideoFileClip(self.filepath)
                if target_fmt in ["mp3", "wav", "flac", "ogg", "aac"]:
                    save_path = os.path.join(out_dir, f"{base_name}.{target_fmt}")
                    clip.audio.write_audiofile(save_path)
                elif target_fmt == "gif":
                    save_path = os.path.join(out_dir, f"{base_name}.gif")
                    clip.write_gif(save_path, fps=10)
                else:
                    save_path = os.path.join(out_dir, f"{base_name}.{target_fmt}")
                    clip.write_videofile(save_path, codec=selected_codec, audio_codec="aac")
                clip.close()

            elif self.file_type == "audio":
                audio_clip = AudioFileClip(self.filepath)
                save_path = os.path.join(out_dir, f"{base_name}.{target_fmt}")
                audio_clip.write_audiofile(save_path)
                audio_clip.close()

            elif self.file_type == "presentation":
                ext = os.path.splitext(self.filepath)[1].lower()
                if ext == ".pdf":
                    doc = pymupdf.open(self.filepath)
                    if target_fmt in ["png", "jpg"]:
                        out_folder = os.path.join(out_dir, f"{base_name}_slides")
                        os.makedirs(out_folder, exist_ok=True)
                        for i, page in enumerate(doc):
                            page.get_pixmap().save(os.path.join(out_folder, f"slide_{i+1}.{target_fmt}"))
                    elif target_fmt == "txt":
                        save_path = os.path.join(out_dir, f"{base_name}.txt")
                        with open(save_path, "w", encoding="utf-8") as f:
                            for page in doc:
                                f.write(page.get_text() + "\n--- PAGE BREAK ---\n")
                    doc.close()
                elif ext == ".pptx":
                    if target_fmt == "txt" and Presentation is not None:
                        prs = Presentation(self.filepath)
                        save_path = os.path.join(out_dir, f"{base_name}.txt")
                        with open(save_path, "w", encoding="utf-8") as f:
                            for i, slide in enumerate(prs.slides):
                                f.write(f"--- SLIDE {i+1} ---\n")
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        f.write(shape.text + "\n")

            self.after(0, lambda: self.lbl_status_conv.configure(text=self._("success_saved")))
            self.after(0, lambda: self.progress_conv.set(1.0))
            self.after(0, lambda: messagebox.showinfo(self._("done"), self._("file_converted")))
        except Exception as e:
            self.after(0, lambda: messagebox.showerror(self._("error"), f"{self._('fail')} {str(e)}"))
        finally:
            self.after(0, lambda: self.btn_convert.configure(state="normal"))

    def choose_download_folder(self):
        folder = filedialog.askdirectory()
        if folder:
            self.download_folder = folder
            self.lbl_folder_path.configure(text=f"{self._('folder_lbl')} {folder}")
            self.save_settings()

    def start_download_thread(self):
        url = self.url_entry.get().strip()
        if not url:
            messagebox.showwarning(self._("warning"), self._("enter_url"))
            return

        self.btn_download.configure(state="disabled")
        self.progress_dl.set(0.1)
        self.lbl_status_dl.configure(text=self._("connecting"))
        threading.Thread(target=self.run_download, args=(url,), daemon=True).start()

    def run_download(self, url):
        selected_option = self.dl_type_combo.get()
        selected_codec = self.dl_codec_combo.get()
        target_ext = selected_option.split()[0].lower()
        is_audio = target_ext in ["mp3", "wav", "flac", "m4a"]

        ydl_opts = {"outtmpl": os.path.join(self.download_folder, "%(title)s.%(ext)s"), "quiet": True, "no_warnings": True}

        if is_audio:
            ydl_opts.update({
                "format": "bestaudio/best",
                "postprocessors": [{"key": "FFmpegExtractAudio", "preferredcodec": target_ext, "preferredquality": "192"}],
            })
        else:
            vcodec_filter = "[vcodec^=avc1]" if "H.264" in selected_codec else "[vcodec^=hev1]" if "H.265" in selected_codec else "[vcodec^=vp09]" if "VP9" in selected_codec else "[vcodec^=av01]" if "AV1" in selected_codec else ""
            ydl_opts.update({
                "format": f"bestvideo{vcodec_filter}+bestaudio/best{vcodec_filter}/best",
                "merge_output_format": target_ext,
            })

        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                self.after(0, lambda: self.lbl_status_dl.configure(text=f"{self._('downloading')} ({target_ext.upper()})..."))
                self.after(0, lambda: self.progress_dl.set(0.6))
                ydl.download([url])

            self.after(0, lambda: self.progress_dl.set(1.0))
            self.after(0, lambda: self.lbl_status_dl.configure(text=self._("dl_complete")))
            self.after(0, lambda: messagebox.showinfo(self._("success"), self._("media_downloaded")))
        except Exception as e:
            self.after(0, lambda: self.progress_dl.set(0))
            self.after(0, lambda: messagebox.showerror(self._("dl_error"), f"{self._('dl_failed')}\n{str(e)}"))
        finally:
            self.after(0, lambda: self.btn_download.configure(state="normal"))

if __name__ == "__main__":
    app = luvkittysApp()
    app.mainloop()